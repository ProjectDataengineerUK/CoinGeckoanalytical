"""
Run:  pip install python-pptx && python3 docs/presentation/coingeckoanalytical_pitch.py
Output: docs/presentation/coingeckoanalytical.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x1F, 0x3B, 0x6E)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF3, 0xF4, 0xF6)
DGRAY  = RGBColor(0x37, 0x41, 0x51)
GREEN  = RGBColor(0x05, 0x96, 0x69)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
TEAL   = RGBColor(0x0F, 0x76, 0x6E)

W = Inches(13.33)
H = Inches(7.5)


def prs() -> Presentation:
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p


def blank(p: Presentation):
    return p.slides.add_slide(p.slide_layouts[6])  # completely blank


def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height,
        fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def label(slide, text, left, top, width, height,
          font_size=Pt(16), bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def title_bar(slide, title: str, subtitle: str = ""):
    box(slide, 0, 0, W, Inches(1.1), fill_color=NAVY)
    label(slide, title, Inches(0.4), Inches(0.12), Inches(12), Inches(0.55),
          font_size=Pt(26), bold=True, color=WHITE)
    if subtitle:
        label(slide, subtitle, Inches(0.4), Inches(0.66), Inches(12), Inches(0.38),
              font_size=Pt(14), color=AMBER)


def bullets(slide, items: list[str], left, top, width, height,
            size=Pt(15), color=DGRAY, indent=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size  = size
        p.font.color.rgb = color
        if indent:
            p.level = 1


def accent_rect(slide, left, top, width, height, color: RGBColor, text: str,
                font_size=Pt(13), text_color=WHITE):
    box(slide, left, top, width, height, fill_color=color)
    label(slide, text, left + Inches(0.1), top + Inches(0.08),
          width - Inches(0.2), height - Inches(0.1),
          font_size=font_size, bold=True, color=text_color,
          align=PP_ALIGN.CENTER)


# ===========================================================================
# Slides
# ===========================================================================

def slide_cover(p: Presentation):
    s = blank(p)
    bg(s, NAVY)
    box(s, 0, Inches(2.8), W, Inches(0.06), fill_color=AMBER)
    label(s, "CoinGeckoAnalytical",
          Inches(1), Inches(1.2), Inches(11), Inches(1.4),
          font_size=Pt(52), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    label(s, "Crypto Market Intelligence Platform",
          Inches(1), Inches(2.9), Inches(11), Inches(0.8),
          font_size=Pt(24), color=AMBER, align=PP_ALIGN.CENTER)
    label(s, "Powered by Databricks  ·  MLOps Level 5  ·  22 Production Jobs",
          Inches(1), Inches(3.8), Inches(11), Inches(0.6),
          font_size=Pt(15), color=RGBColor(0xA0, 0xB4, 0xD6), align=PP_ALIGN.CENTER)
    label(s, "cga-analytics  ·  cga-admin  ·  Genie BI  ·  Multi-Agent Copilot",
          Inches(1), Inches(4.5), Inches(11), Inches(0.6),
          font_size=Pt(13), color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.CENTER)


def slide_overview(p: Presentation):
    s = blank(p)
    bg(s, LGRAY)
    title_bar(s, "Product Overview", "What it is · who it's for · what problem it solves")

    label(s, "CoinGeckoAnalytical aggregates crypto market data from four external sources, "
             "runs a full medallion data pipeline, and surfaces structured analytics + AI narrative "
             "insights through two dedicated Databricks Apps.",
          Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.9),
          font_size=Pt(15), color=DGRAY)

    cols = [
        (NAVY,   "Analysts & Traders",    "Real-time market\nrankings, DeFi TVL,\nmacro overlays"),
        (GREEN,  "Ops Teams",             "Pipeline health,\nSentinela alerts,\ncost telemetry"),
        (PURPLE, "Platform Engineers",    "Governance, RLS,\naccess management,\naudit review"),
        (TEAL,   "Institutional Teams",   "Regime classification,\nanomaly detection,\nML scoring"),
    ]
    for i, (c, title, desc) in enumerate(cols):
        x = Inches(0.3 + i * 3.18)
        box(s, x, Inches(2.3), Inches(3.0), Inches(3.8), fill_color=c)
        label(s, title, x + Inches(0.1), Inches(2.42), Inches(2.8), Inches(0.5),
              font_size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        label(s, desc, x + Inches(0.1), Inches(2.98), Inches(2.8), Inches(2.8),
              font_size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)


def slide_architecture(p: Presentation):
    s = blank(p)
    bg(s, LGRAY)
    title_bar(s, "High-Level Architecture", "Three planes — data · AI · product")

    layers = [
        (RGBColor(0x6B, 0x72, 0x80), "DATA SOURCES",  "CoinGecko · DefiLlama · GitHub · FRED"),
        (RGBColor(0x92, 0x40, 0x0E), "BRONZE",         "Raw ingestion · MERGE INTO · Delta Lake"),
        (RGBColor(0x64, 0x74, 0x8B), "SILVER",         "Normalised · enriched · validated"),
        (RGBColor(0xD9, 0x77, 0x06), "GOLD",            "Analytical models · Genie metric views"),
        (RGBColor(0x1D, 0x4E, 0xD8), "UNITY CATALOG",  "Governance · RLS · column masking · lineage"),
        (PURPLE,                      "AI LAYER",        "MLflow Registry · Mosaic AI · Copilot Orchestrator"),
        (GREEN,                       "PRODUCT SURFACE", "cga-analytics (Genie + Charts + Copilot)  ·  cga-admin (Sentinela + Ops)"),
    ]
    for i, (c, name, desc) in enumerate(layers):
        y = Inches(1.2 + i * 0.84)
        box(s, Inches(0.4), y, Inches(2.2), Inches(0.7), fill_color=c)
        label(s, name, Inches(0.5), y + Inches(0.08), Inches(2.0), Inches(0.55),
              font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        label(s, desc, Inches(2.9), y + Inches(0.1), Inches(9.8), Inches(0.55),
              font_size=Pt(13), color=DGRAY)
        if i < len(layers) - 1:
            label(s, "▼", Inches(1.3), y + Inches(0.72), Inches(0.5), Inches(0.3),
                  font_size=Pt(10), color=DGRAY, align=PP_ALIGN.CENTER)


def slide_pipeline(p: Presentation):
    s = blank(p)
    bg(s, LGRAY)
    title_bar(s, "Data Pipeline — 22 Jobs", "Bronze → Silver → Gold · 4 external sources · Databricks Asset Bundles")

    phases = [
        (RGBColor(0x92, 0x40, 0x0E), "BRONZE\nINGESTION",
         ["market_source_ingestion_job", "defillama_ingestion_job",
          "github_activity_ingestion_job", "fred_macro_ingestion_job",
          "bronze_market_table_migration_job", "bronze_enrichment_migration_job"]),
        (RGBColor(0x64, 0x74, 0x8B), "SILVER\nPIPELINES",
         ["silver_market_pipeline_job", "silver_enrichment_pipeline_job",
          "silver_market_table_migration_job", "silver_enrichment_migration_job"]),
        (RGBColor(0xD9, 0x77, 0x06), "GOLD / OPS",
         ["ops_readiness_refresh_job", "ops_usage_ingestion_job",
          "ops_bundle_run_ingestion_job", "ops_sentinela_alert_ingestion_job",
          "sentinela_evaluation_job"]),
        (PURPLE, "MLOPS",
         ["feature_engineering_job", "train_market_model_job",
          "score_market_assets_job", "model_drift_monitoring_job"]),
        (RGBColor(0x1D, 0x4E, 0xD8), "GOVERNANCE",
         ["uc_grants_job", "rls_migration_job", "sentinela_evaluation_job*"]),
    ]
    for i, (c, name, jobs) in enumerate(phases):
        x = Inches(0.25 + i * 2.6)
        box(s, x, Inches(1.25), Inches(2.45), Inches(5.7), fill_color=c)
        label(s, name, x + Inches(0.1), Inches(1.32), Inches(2.25), Inches(0.65),
              font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, job in enumerate(jobs):
            box(s, x + Inches(0.1), Inches(2.05 + j * 0.77), Inches(2.25), Inches(0.65),
                fill_color=WHITE)
            label(s, job.replace("_job", "").replace("_", " "),
                  x + Inches(0.15), Inches(2.1 + j * 0.77), Inches(2.15), Inches(0.55),
                  font_size=Pt(9), color=c, bold=True)


def slide_analytics_app(p: Presentation):
    s = blank(p)
    bg(s, LGRAY)
    title_bar(s, "cga-analytics — User App", "Genie BI controller · Dynamic charts · Multi-agent copilot")

    panels = [
        (NAVY,   "Genie Panel",
         ["Natural language query input", "genie_client.ask_genie()", "Returns generated SQL",
          "Drives shared active_dataset state", "Whole dashboard re-renders on each query"]),
        (RGBColor(0x0E, 0x7A, 0x4A), "Chart Dashboard",
         ["Subscribes to active_dataset state", "Market rankings (top 50)", "DeFi TVL overlay",
          "Macro indicators (Fed Rate, CPI, M2)", "Price change heatmap"]),
        (PURPLE, "Copilot Panel",
         ["Tier routing: light / standard / complex", "light → direct Mosaic AI",
          "complex → 3-agent orchestrator", "Answer: narrative + provenance + freshness",
          "Cost telemetry on every response"]),
    ]
    for i, (c, title, pts) in enumerate(panels):
        x = Inches(0.3 + i * 4.3)
        box(s, x, Inches(1.25), Inches(4.1), Inches(5.7), fill_color=c)
        label(s, title, x + Inches(0.1), Inches(1.35), Inches(3.9), Inches(0.55),
              font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, pt in enumerate(pts):
            label(s, f"• {pt}", x + Inches(0.15), Inches(2.05 + j * 0.85), Inches(3.8), Inches(0.75),
                  font_size=Pt(12), color=WHITE)


def slide_admin_app(p: Presentation):
    s = blank(p)
    bg(s, LGRAY)
    title_bar(s, "cga-admin — Ops App", "Sentinela monitoring · cost telemetry · access management")

    pages = [
        (RGBColor(0xDC, 0x26, 0x26), "Sentinela Monitor",
         "Live alerts · CRITICAL/WARNING/OK\nsentinela_evaluation_job every 15 min"),
        (TEAL,                         "Pipeline Health",
         "Job run history · failure rates\nbundle_run_ingestion observability"),
        (RGBColor(0xD9, 0x77, 0x06),  "Cost Monitor",
         "Token usage per tier · cost_usd\nColumn masking for non-ops users"),
        (RGBColor(0x1D, 0x4E, 0xD8),  "Access Management",
         "UC grants · tenant controls\nRLS row filters per user group"),
        (PURPLE,                       "Audit Review",
         "ops_usage_events · query log\nFull audit trail per user_id (hashed)"),
    ]
    for i, (c, title, desc) in enumerate(pages):
        x = Inches(0.25 + (i % 3) * 4.28)
        y = Inches(1.3 if i < 3 else 4.1)
        box(s, x, y, Inches(4.0), Inches(2.55), fill_color=c)
        label(s, title, x + Inches(0.15), y + Inches(0.12), Inches(3.7), Inches(0.5),
              font_size=Pt(15), bold=True, color=WHITE)
        label(s, desc, x + Inches(0.15), y + Inches(0.75), Inches(3.7), Inches(1.6),
              font_size=Pt(12), color=WHITE)


def slide_ai_copilot(p: Presentation):
    s = blank(p)
    bg(s, LGRAY)
    title_bar(s, "AI / Copilot Architecture", "Coded Python orchestrator · tier routing · provenance-aware answers")

    box(s, Inches(0.3), Inches(1.25), Inches(12.7), Inches(5.7), fill_color=WHITE,
        line_color=NAVY, line_width=Pt(1.5))

    label(s, "User Query", Inches(0.5), Inches(1.4), Inches(2.0), Inches(0.5),
          font_size=Pt(13), bold=True, color=NAVY)
    accent_rect(s, Inches(0.5), Inches(2.0), Inches(2.0), Inches(0.8), TEAL,
                "model_tier_router", font_size=Pt(11))
    label(s, "↓", Inches(1.3), Inches(2.85), Inches(0.5), Inches(0.4),
          font_size=Pt(16), color=DGRAY, align=PP_ALIGN.CENTER)

    tiers = [
        (Inches(0.4), RGBColor(0x05, 0x96, 0x69), "LIGHT\nTier",  "Direct\nMosaic AI"),
        (Inches(3.1), RGBColor(0xD9, 0x77, 0x06), "STANDARD\nTier", "Mosaic AI\n+ context"),
        (Inches(5.8), PURPLE,                      "COMPLEX\nTier", "Orchestrator\n3 agents"),
    ]
    for x, c, t, d in tiers:
        accent_rect(s, x, Inches(3.35), Inches(2.3), Inches(0.75), c, t, font_size=Pt(11))
        label(s, d, x + Inches(0.1), Inches(4.2), Inches(2.1), Inches(0.6),
              font_size=Pt(11), color=c, align=PP_ALIGN.CENTER)

    agents = ["MarketAgent\ngold_market_*", "MacroAgent\ngold_macro_*",
              "DefiAgent\ngold_defi_*", "SynthAgent\n(synthesis)"]
    for i, ag in enumerate(agents):
        x = Inches(6.0 + i * 1.7)
        accent_rect(s, x, Inches(4.95), Inches(1.55), Inches(1.7), NAVY, ag, font_size=Pt(10))

    label(s, "Answer + provenance + freshness badge + cost_usd",
          Inches(0.5), Inches(6.75), Inches(10), Inches(0.45),
          font_size=Pt(13), bold=True, color=GREEN)


def slide_mlops(p: Presentation):
    s = blank(p)
    bg(s, LGRAY)
    title_bar(s, "MLOps Pipeline", "Level 5 — cross-val · PSI drift · MLflow registry · batch scoring")

    steps = [
        (RGBColor(0x64, 0x74, 0x8B), "Feature\nEngineering",
         "Silver → feature_store\nmomentum / dominance\ncross-asset correlation"),
        (PURPLE, "Train",
         "Regime Classifier\n(bull/bear/neutral)\nAnomaly Detector\ncross_val_score ≥ 0.60"),
        (RGBColor(0xD9, 0x77, 0x06), "MLflow\nRegistry",
         "Experiment tracking\nModel versioning\nChampion promotion\nvalidation gate"),
        (GREEN,  "Score",
         "Batch scoring\ngold_ml_scores\nregime + anomaly\nper asset"),
        (RGBColor(0xDC, 0x26, 0x26), "Drift\nMonitor",
         "PSI > 0.2 → alert\n24h vs 7d window\nops_sentinela_alerts\ndaily cron"),
    ]
    for i, (c, title, desc) in enumerate(steps):
        x = Inches(0.3 + i * 2.58)
        box(s, x, Inches(1.3), Inches(2.4), Inches(5.7), fill_color=c)
        label(s, title, x + Inches(0.1), Inches(1.4), Inches(2.2), Inches(0.75),
              font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        label(s, desc, x + Inches(0.1), Inches(2.25), Inches(2.2), Inches(4.4),
              font_size=Pt(12), color=WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            label(s, "→", Inches(0.3 + (i + 1) * 2.58 - 0.22), Inches(3.8), Inches(0.3), Inches(0.5),
                  font_size=Pt(20), color=DGRAY, align=PP_ALIGN.CENTER)


def slide_governance(p: Presentation):
    s = blank(p)
    bg(s, LGRAY)
    title_bar(s, "Governance & Compliance", "Unity Catalog · RLS · column masking · uc_grants automation")

    items = [
        (RGBColor(0x1D, 0x4E, 0xD8), "Unity Catalog",
         ["3 catalogs: cgadev / cgastaging / cgaprod",
          "All assets governed under UC",
          "Automated lineage via System Tables",
          "Schema-level GRANT enforcement"]),
        (RGBColor(0xDC, 0x26, 0x26), "Row-Level Security",
         ["UC row filter functions on ops tables",
          "ops_usage_events: platform_ops only",
          "ops_sentinela_alerts: ops/governance only",
          "rls_migration_job applies via CI"]),
        (PURPLE, "Column Masking",
         ["cost_usd masked for non-ops users",
          "CREATE OR REPLACE FUNCTION masks to NULL",
          "ALTER TABLE SET COLUMN MASK",
          "Triggered via confirm_uc_grants gate"]),
        (TEAL,   "uc_grants Automation",
         ["unity_catalog_foundation.sql",
          "All GRANT statements version-controlled",
          "CI-triggerable: confirm_uc_grants=true",
          "No manual GRANT commands ever"]),
    ]
    for i, (c, title, pts) in enumerate(items):
        x = Inches(0.3 + (i % 2) * 6.4)
        y = Inches(1.3 if i < 2 else 4.1)
        box(s, x, y, Inches(6.1), Inches(2.55), fill_color=c)
        label(s, title, x + Inches(0.15), y + Inches(0.1), Inches(5.8), Inches(0.5),
              font_size=Pt(15), bold=True, color=WHITE)
        for j, pt in enumerate(pts):
            label(s, f"• {pt}", x + Inches(0.2), y + Inches(0.7 + j * 0.42), Inches(5.7), Inches(0.4),
                  font_size=Pt(12), color=WHITE)


def slide_cicd(p: Presentation):
    s = blank(p)
    bg(s, LGRAY)
    title_bar(s, "CI/CD & Deploy", "GitHub Actions · Databricks Asset Bundles · workflow_dispatch gates")

    jobs_ci = [
        (DGRAY,  "lint",          "python -m compileall\nall sources"),
        (TEAL,   "contract",      "346 tests\nbundle validation\nchain validators"),
        (NAVY,   "deploy",        "bundle deploy -t dev\n20 pipeline jobs\nlive SQL validation"),
        (GREEN,  "deploy_apps",   "cga-analytics start\ncga-admin start\napps verify"),
        (RGBColor(0x1D, 0x4E, 0xD8), "uc_grants", "uc_grants_job\nrls_migration_job\nGRANT execution"),
        (PURPLE, "train_models",  "train_market_model\nMLflow promotion\nchampion gate"),
    ]
    prev_x = None
    for i, (c, name, desc) in enumerate(jobs_ci):
        x = Inches(0.3 + i * 2.17)
        box(s, x, Inches(1.3), Inches(2.0), Inches(4.5), fill_color=c)
        label(s, name, x + Inches(0.1), Inches(1.4), Inches(1.8), Inches(0.55),
              font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        label(s, desc, x + Inches(0.1), Inches(2.05), Inches(1.8), Inches(3.5),
              font_size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)
        if prev_x is not None and i < 3:
            label(s, "→", x - Inches(0.18), Inches(3.2), Inches(0.25), Inches(0.4),
                  font_size=Pt(18), color=DGRAY, align=PP_ALIGN.CENTER)
        prev_x = x

    label(s, "setup-cli@main  ·  confirm_deploy  ·  confirm_apps_deploy  ·  confirm_uc_grants  ·  confirm_train",
          Inches(0.4), Inches(6.1), Inches(12.4), Inches(0.5),
          font_size=Pt(12), color=NAVY, align=PP_ALIGN.CENTER)
    label(s, "deploy_apps, uc_grants, train_models run AFTER deploy completes (serialised — no lock conflict)",
          Inches(0.4), Inches(6.6), Inches(12.4), Inches(0.5),
          font_size=Pt(11), color=DGRAY, align=PP_ALIGN.CENTER)


def slide_maturity(p: Presentation):
    s = blank(p)
    bg(s, LGRAY)
    title_bar(s, "Maturity Level 5 — DataOps · MLOps · LLMOps", "All three disciplines at production grade")

    cols = [
        (TEAL,   "DataOps L5",
         ["MERGE INTO — idempotent ingestion", "Post-write table count (Serverless safe)",
          "Data contracts: 7 JSON Schema files", "Great Expectations-style assertions",
          "live_sql_validation.py against real warehouse",
          "ops_readiness_refresh — Gold views scheduled",
          "Sentinela evaluation every 15 min",
          "Full lineage in Unity Catalog System Tables"]),
        (PURPLE, "MLOps L5",
         ["Cross-validation accuracy (cv=5)", "Validation gate: cv_accuracy ≥ 0.60",
          "PSI drift monitoring (threshold 0.2)", "MLflow Model Registry + champion promotion",
          "Batch scoring → gold_ml_scores", "Weekly retraining cron (Mondays 2am)",
          "Feature store: momentum + dominance features",
          "model_drift_monitoring_job daily"]),
        (NAVY,   "LLMOps L5",
         ["Versioned prompts: backend/prompts/v1.yaml", "sanitize_user_input() injection guard",
          "Per-tier token budget cap", "user_id hashed (sha256) in telemetry",
          "Golden eval set: backend/eval/golden_eval.json",
          "cost_usd on every AI response", "Tier routing: light/standard/complex",
          "Exception logging (not swallowed)"]),
    ]
    for i, (c, title, pts) in enumerate(cols):
        x = Inches(0.25 + i * 4.35)
        box(s, x, Inches(1.25), Inches(4.2), Inches(6.0), fill_color=c)
        label(s, title, x + Inches(0.1), Inches(1.35), Inches(4.0), Inches(0.55),
              font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, pt in enumerate(pts):
            label(s, f"• {pt}", x + Inches(0.15), Inches(2.0 + j * 0.63), Inches(3.9), Inches(0.58),
                  font_size=Pt(11), color=WHITE)


def slide_security(p: Presentation):
    s = blank(p)
    bg(s, LGRAY)
    title_bar(s, "Security Posture", "SQL injection · thread safety · prompt injection · data masking")

    fixes = [
        (RGBColor(0xDC, 0x26, 0x26), "SQL Injection",
         "_SAFE_ASSET_ID_RE regex\nvalidates asset_ids before\nIN clause interpolation"),
        (TEAL, "Thread Safety",
         "_TOKEN_CACHE_LOCK\n(threading.Lock) wraps\nall token cache R/W"),
        (PURPLE, "Prompt Injection",
         "sanitize_user_input()\nstrips control chars +\nrole-override patterns"),
        (NAVY, "Data Masking",
         "cost_usd masked to NULL\nfor non-ops users via\nUC column mask function"),
        (GREEN, "User Privacy",
         "user_id hashed sha256\nbefore storage in\nops_usage_events"),
        (RGBColor(0xD9, 0x77, 0x06), "Auth",
         "OAuth M2M / service\nprincipal in CI — no\nuser PAT in automation"),
    ]
    for i, (c, title, desc) in enumerate(fixes):
        x = Inches(0.3 + (i % 3) * 4.3)
        y = Inches(1.3 if i < 3 else 4.15)
        box(s, x, y, Inches(4.1), Inches(2.6), fill_color=c)
        label(s, title, x + Inches(0.15), y + Inches(0.12), Inches(3.8), Inches(0.5),
              font_size=Pt(15), bold=True, color=WHITE)
        label(s, desc, x + Inches(0.15), y + Inches(0.72), Inches(3.8), Inches(1.7),
              font_size=Pt(13), color=WHITE)


def slide_metrics(p: Presentation):
    s = blank(p)
    bg(s, LGRAY)
    title_bar(s, "Project Metrics & Status", "All 15 blocks DONE · 346 tests · 22 jobs · CI green")

    metrics = [
        (NAVY,   "22", "Production Jobs"),
        (GREEN,  "346", "Tests Passing"),
        (PURPLE, "15+", "Project Blocks DONE"),
        (RGBColor(0xDC, 0x26, 0x26), "4", "Data Sources"),
        (TEAL,   "2", "Databricks Apps"),
        (RGBColor(0xD9, 0x77, 0x06), "5", "Maturity Level"),
    ]
    for i, (c, num, label_text) in enumerate(metrics):
        x = Inches(0.4 + (i % 3) * 4.27)
        y = Inches(1.3 if i < 3 else 4.0)
        box(s, x, y, Inches(4.0), Inches(2.5), fill_color=c)
        label(s, num, x, y + Inches(0.1), Inches(4.0), Inches(1.4),
              font_size=Pt(56), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        label(s, label_text, x, y + Inches(1.6), Inches(4.0), Inches(0.7),
              font_size=Pt(16), color=WHITE, align=PP_ALIGN.CENTER)


def slide_nextup(p: Presentation):
    s = blank(p)
    bg(s, NAVY)
    box(s, 0, Inches(3.55), W, Inches(0.06), fill_color=AMBER)
    label(s, "Recommended Next Steps",
          Inches(0.5), Inches(0.6), Inches(12), Inches(0.9),
          font_size=Pt(32), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    next_steps = [
        ("🚀 Staging promotion",    "Run full pipeline on cgastaging catalog with real CoinGecko API key"),
        ("🔑 OAuth M2M setup",       "Configure service principal client_id/secret in GitHub Secrets"),
        ("📊 Genie space config",    "Create AI/BI Genie space linked to Gold views in cga-analytics"),
        ("🔔 Notification hooks",    "Wire Sentinela CRITICAL alerts to Slack / PagerDuty webhook"),
        ("🌍 Multi-source expansion","Add Binance, CMC, or additional on-chain data providers"),
        ("📱 Mobile / embed",        "Embed cga-analytics in iframe or export Genie as public dashboard"),
    ]
    for i, (title, desc) in enumerate(next_steps):
        y = Inches(1.6 + i * 0.9)
        x_t = Inches(0.6)
        x_d = Inches(4.2)
        label(s, title, x_t, y, Inches(3.4), Inches(0.75),
              font_size=Pt(14), bold=True, color=AMBER)
        label(s, desc, x_d, y, Inches(8.5), Inches(0.75),
              font_size=Pt(13), color=RGBColor(0xD1, 0xD5, 0xDB))


# ===========================================================================
# Main
# ===========================================================================

def main() -> None:
    p = prs()
    slide_cover(p)
    slide_overview(p)
    slide_architecture(p)
    slide_pipeline(p)
    slide_analytics_app(p)
    slide_admin_app(p)
    slide_ai_copilot(p)
    slide_mlops(p)
    slide_governance(p)
    slide_cicd(p)
    slide_maturity(p)
    slide_security(p)
    slide_metrics(p)
    slide_nextup(p)

    out = Path(__file__).resolve().parent / "coingeckoanalytical.pptx"
    p.save(str(out))
    print(f"Saved: {out}  ({len(p.slides)} slides)")


if __name__ == "__main__":
    main()
